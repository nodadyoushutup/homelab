"""Structured chunking for Office Open XML (DOCX, PPTX) and ODT (text extraction + batching)."""
from __future__ import annotations

import io
import logging
import os
import zipfile
from typing import Any

from lxml import etree

from chunks.structured import StructuredChunk, _split_oversized
from chunks.tabular import md_table

log = logging.getLogger(__name__)


def _env_int(name: str, default: int) -> int:
    try:
        return int(os.getenv(name, str(default)).strip())
    except ValueError:
        return default


def office_docx_profile() -> str:
    p = max(1, _env_int("RAG_OFFICE_DOCX_PARAS_PER_CHUNK", 12))
    return f"docx|{p}"


def office_pptx_profile() -> str:
    n = max(1, _env_int("RAG_OFFICE_PPTX_SLIDES_PER_CHUNK", 1))
    return f"pptx|{n}"


def office_odt_profile() -> str:
    p = max(1, _env_int("RAG_OFFICE_ODT_PARAS_PER_CHUNK", 12))
    return f"odt|{p}"


def _batch_segments(
    rel_path: str,
    *,
    language: str,
    chunk_kind: str,
    segments: list[str],
    profile: str,
    max_chars: int,
    overlap: int,
) -> list[StructuredChunk]:
    out: list[StructuredChunk] = []
    batch: list[str] = []
    blen = 0
    part = 0

    def flush() -> None:
        nonlocal batch, blen, part
        if not batch:
            return
        body = "\n\n".join(batch).strip()
        part += 1
        hdr_kw: dict[str, Any] = {
            "office_part": "batched",
            "office_index": part,
            "office_ingest_profile": profile,
        }
        out.extend(
            _split_oversized(
                rel_path,
                language,
                chunk_kind,
                f"part{part}",
                body + "\n",
                part,
                max_chars=max_chars,
                overlap=overlap,
                header_kwargs=hdr_kw,
            )
        )
        batch = []
        blen = 0

    for seg in segments:
        s = seg.strip()
        if not s:
            continue
        add_len = len(s) + 2
        if batch and blen + add_len > max_chars:
            flush()
        batch.append(s)
        blen += add_len
    flush()
    return out


def build_docx_chunks(rel_path: str, raw: bytes, *, max_chars: int, overlap: int) -> list[StructuredChunk]:
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        log.warning("python-docx missing; docx ingest disabled for %s", rel_path)
        return []
    profile = office_docx_profile()
    try:
        doc = Document(io.BytesIO(raw))
    except Exception as exc:
        log.info("docx open failed for %s: %s", rel_path, exc)
        return []

    segments: list[str] = []
    body = doc.element.body
    for child in body.iterchildren():
        tag = child.tag
        if tag == qn("w:p"):
            p = Paragraph(child, doc)
            t = p.text.strip()
            if t:
                try:
                    st = p.style.name if p.style is not None else ""
                except (AttributeError, ValueError):
                    st = ""
                if st and "Heading" in st:
                    try:
                        lvl = int(st.replace("Heading", "").strip() or "1")
                    except ValueError:
                        lvl = 1
                    segments.append(f"{'#' * max(1, min(lvl, 6))} {t}")
                else:
                    segments.append(t)
        elif tag == qn("w:tbl"):
            tbl = Table(child, doc)
            rows: list[list[str]] = []
            for row in tbl.rows:
                rows.append([c.text.strip().replace("\n", " ") for c in row.cells])
            if not rows:
                continue
            header = rows[0]
            body_rows = rows[1:] if len(rows) > 1 else []
            md = md_table(header, body_rows)
            if md:
                segments.append(md)

    return _batch_segments(
        rel_path,
        language="docx",
        chunk_kind="docx_block",
        segments=segments,
        profile=profile,
        max_chars=max_chars,
        overlap=overlap,
    )


def build_pptx_chunks(rel_path: str, raw: bytes, *, max_chars: int, overlap: int) -> list[StructuredChunk]:
    try:
        from pptx import Presentation
    except ImportError:
        log.warning("python-pptx missing; pptx ingest disabled for %s", rel_path)
        return []
    profile = office_pptx_profile()
    slides_per = max(1, _env_int("RAG_OFFICE_PPTX_SLIDES_PER_CHUNK", 1))
    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as exc:
        log.info("pptx open failed for %s: %s", rel_path, exc)
        return []

    slide_texts: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            tx = (shape.text or "").strip()
            if tx:
                parts.append(tx)
        slide_texts.append("\n".join(parts).strip())

    segments: list[str] = []
    buf: list[str] = []
    for i, st in enumerate(slide_texts):
        if not st:
            continue
        labeled = f"--- slide {i + 1} ---\n{st}"
        buf.append(labeled)
        if len(buf) >= slides_per:
            segments.append("\n\n".join(buf))
            buf = []
    if buf:
        segments.append("\n\n".join(buf))

    return _batch_segments(
        rel_path,
        language="pptx",
        chunk_kind="pptx_slides",
        segments=segments,
        profile=profile,
        max_chars=max_chars,
        overlap=overlap,
    )


_OFFICE_NS = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
_TEXT_NS = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
_TABLE_NS = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"


def build_odt_chunks(rel_path: str, raw: bytes, *, max_chars: int, overlap: int) -> list[StructuredChunk]:
    profile = office_odt_profile()
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            xml = zf.read("content.xml")
    except (KeyError, zipfile.BadZipFile, OSError) as exc:
        log.info("odt zip read failed for %s: %s", rel_path, exc)
        return []

    try:
        root = etree.fromstring(xml)
    except etree.XMLSyntaxError as exc:
        log.info("odt content.xml parse failed for %s: %s", rel_path, exc)
        return []

    body_text = root.find(f".//{{{_OFFICE_NS}}}text")
    if body_text is None:
        log.info("odt: missing office:text in %s", rel_path)
        return []

    segments: list[str] = []

    def text_el(el: Any) -> str:
        return "".join(el.itertext()).strip()

    for el in body_text:
        tag = el.tag
        if tag == f"{{{_TEXT_NS}}}p":
            t = text_el(el)
            if t:
                segments.append(t)
        elif tag == f"{{{_TEXT_NS}}}h":
            t = text_el(el)
            if t:
                lvl = 1
                ol = el.get(f"{{{_TEXT_NS}}}outline-level")
                if ol and str(ol).isdigit():
                    lvl = max(1, min(6, int(ol)))
                segments.append(f"{'#' * lvl} {t}")
        elif tag == f"{{{_TABLE_NS}}}table":
            rows_out: list[list[str]] = []
            for row in el.iterchildren():
                if row.tag != f"{{{_TABLE_NS}}}table-row":
                    continue
                cells: list[str] = []
                for cell in row.iterchildren():
                    if cell.tag == f"{{{_TABLE_NS}}}table-cell":
                        cells.append(" ".join(cell.itertext()).strip().replace("\n", " "))
                if cells:
                    rows_out.append(cells)
            if rows_out:
                header = rows_out[0]
                body_rows = rows_out[1:] if len(rows_out) > 1 else []
                md = md_table(header, body_rows)
                if md:
                    segments.append(md)

    return _batch_segments(
        rel_path,
        language="odt",
        chunk_kind="odt_block",
        segments=segments,
        profile=profile,
        max_chars=max_chars,
        overlap=overlap,
    )
